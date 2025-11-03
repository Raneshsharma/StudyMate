"""
Streamlit RAG Study Assistant — Stable Version (No Chroma Warning)
- Uses PersistentClient (new Chroma API)
- Reads OpenAI API key from Streamlit Secrets or environment
- Clean off-white UI
- Learner profile (Name, Age, Profession, Study Goal)
- File upload (PDF, DOCX, PPTX, TXT)
- RAG chat, summaries, export to .docx

To run:
1. pip install -r requirements.txt
2. Add your OpenAI API key in .streamlit/secrets.toml or environment variable OPENAI_API_KEY
3. streamlit run streamlit_rag_app.py
"""

import streamlit as st
from typing import List, Dict, Any
import os
import tempfile
import uuid
import time
from datetime import datetime
from pathlib import Path
import io
import hashlib

# Document parsing
import pdfplumber
import docx
from pptx import Presentation

# LLM and embeddings
import openai
import chromadb
from chromadb.config import Settings, DEFAULT_TENANT, DEFAULT_DATABASE
from chromadb.utils import embedding_functions

# Export
from docx import Document

# --------------------------- Setup ---------------------------

st.set_page_config(page_title="RAG Study Assistant", layout="wide")

CHROMA_DIR = "./chroma_db"
os.makedirs(CHROMA_DIR, exist_ok=True)

# Read API key from secrets or environment
OPENAI_KEY = st.secrets.get("OPENAI_API_KEY", os.getenv("OPENAI_API_KEY"))
if not OPENAI_KEY:
    st.error("Missing OpenAI API key. Add it to Streamlit Secrets or set environment variable OPENAI_API_KEY.")
else:
    openai.api_key = OPENAI_KEY

# Styling
st.markdown("""
<style>
body { background-color: #f8f6f3; }
[data-testid="stSidebar"] { background-color: #fff; }
.card { background: white; border-radius: 12px; padding: 18px; box-shadow: 0 6px 18px rgba(0,0,0,0.06); }
.h1 { font-size: 26px; font-weight:700; }
.small-muted { color: #6b6b6b; font-size:13px; }
.person-pill { background:#f0f0f0; padding:8px 12px; border-radius:999px; display:inline-block; }
</style>
""", unsafe_allow_html=True)

# --------------------------- Utility Functions ---------------------------

def chunk_text(text: str, chunk_size: int = 1200, overlap: int = 200) -> List[str]:
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start = end - overlap
        if start < 0:
            start = 0
    return chunks

def extract_text(file_bytes: bytes, ext: str) -> str:
    if ext == 'pdf':
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            return "\n".join([p.extract_text() or "" for p in pdf.pages])
    elif ext in ['docx', 'doc']:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            doc = docx.Document(tmp.name)
            text = "\n".join([p.text for p in doc.paragraphs])
        os.unlink(tmp.name)
        return text
    elif ext == 'pptx':
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            prs = Presentation(tmp.name)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
        os.unlink(tmp.name)
        return "\n".join(texts)
    else:
        return file_bytes.decode("utf-8", errors="ignore")

def get_chroma_collection():
    client = chromadb.PersistentClient(
        path=CHROMA_DIR,
        settings=Settings(),
        tenant=DEFAULT_TENANT,
        database=DEFAULT_DATABASE,
    )
    ef = embedding_functions.OpenAIEmbeddingFunction(api_key=openai.api_key, model_name="text-embedding-3-small")
    try:
        collection = client.get_collection(name="documents")
    except Exception:
        collection = client.create_collection(name="documents", embedding_function=ef)
    return client, collection

def add_docs(collection, filename, chunks):
    ids = [f"{filename}_{i}" for i in range(len(chunks))]
    metas = [{"source": filename, "chunk_index": i} for i in range(len(chunks))]
    collection.add(ids=ids, documents=chunks, metadatas=metas)

def query_docs(collection, query: str, k: int = 4):
    if collection.count() == 0:
        return []
    res = collection.query(query_texts=[query], n_results=k)
    docs = res.get("documents", [[]])[0]
    metas = res.get("metadatas", [[]])[0]
    dists = res.get("distances", [[]])[0]
    return [{"text": d, "meta": m, "dist": dist} for d, m, dist in zip(docs, metas, dists)]

def chat(messages):
    try:
        resp = openai.ChatCompletion.create(model="gpt-4o-mini", messages=messages, temperature=0.2, max_tokens=512)
        return resp.choices[0].message.content.strip()
    except Exception as e:
        st.error(f"OpenAI Error: {e}")
        return ""

def export_doc(history, fname="session.docx"):
    doc = Document()
    for h in history:
        doc.add_paragraph(f"[{h['ts']}] {h['role'].upper()}:\n{h['content']}\n")
    doc.save(fname)
    return fname

# --------------------------- Initialize ---------------------------
if 'history' not in st.session_state:
    st.session_state.history = []

client, collection = get_chroma_collection()

# --------------------------- UI ---------------------------
left, right = st.columns([2.5, 1])

with left:
    st.markdown("<div class='card'><h1 class='h1'>RAG Study Assistant</h1><p class='small-muted'>Clean UI · Off-white theme · Persistent Chroma</p></div>", unsafe_allow_html=True)

    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.subheader("Learner Profile")
    pname = st.text_input("Name", value="Aman Verma")
    page = st.number_input("Age", 10, 100, 24)
    profession = st.text_input("Profession", value="MBA Student")
    goal = st.text_input("Study Goal", value="Understand key concepts and make concise notes")
    st.session_state.profile = {"name": pname, "age": page, "profession": profession, "goal": goal}
    st.markdown("</div>", unsafe_allow_html=True)

    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.subheader("Upload Study Files")
    uploaded = st.file_uploader("Upload PDF/DOCX/PPTX/TXT", accept_multiple_files=True)
    if uploaded:
        for f in uploaded:
            ext = f.name.split('.')[-1].lower()
            text = extract_text(f.read(), ext)
            chunks = chunk_text(text)
            add_docs(collection, f.name, chunks)
            st.success(f"Indexed: {f.name} ({len(chunks)} chunks)")
    st.markdown("</div>", unsafe_allow_html=True)

    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.subheader("Ask a Question")
    query = st.text_area("Question", height=120)
    if st.button("Ask"):
        context = query_docs(collection, query, k=4)
        context_text = "\n".join([c['text'] for c in context])
        msgs = [{"role": "system", "content": "You are a helpful study assistant."}, {"role": "user", "content": f"Context:\n{context_text}\nQuestion: {query}"}]
        ans = chat(msgs)
        if ans:
            st.session_state.history.append({"role": "user", "content": query, "ts": datetime.utcnow().isoformat()})
            st.session_state.history.append({"role": "assistant", "content": ans, "ts": datetime.utcnow().isoformat()})
            st.write(ans)
    st.markdown("</div>", unsafe_allow_html=True)

    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.subheader("Summarize & Export")
    if st.button("Summarize All Chats"):
        combined = "\n".join([f"Q: {h['content']}" if h['role']=='user' else f"A: {h['content']}" for h in st.session_state.history])
        prompt = f"Summarize the following conversation into clean study notes in simple language:\n{combined}"
        notes = chat([{"role":"system","content":"Summarizer"},{"role":"user","content":prompt}])
        st.write(notes)
        fname = export_doc([{"role":"notes","content":notes,"ts":datetime.utcnow().isoformat()}], fname=f"study_notes_{int(time.time())}.docx")
        with open(fname, 'rb') as fp:
            st.download_button("Download Notes (.docx)", fp, file_name=fname)
    st.markdown("</div>", unsafe_allow_html=True)

with right:
    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.subheader("Profile")
    p = st.session_state.profile
    st.markdown(f"<div class='person-pill'>{p['name']} · {p['profession']} · Age {p['age']}</div>", unsafe_allow_html=True)
    st.write(p['goal'])
    st.markdown("---")
    st.subheader("Recent Activity")
    if st.session_state.history:
        for h in st.session_state.history[::-1][:5]:
            st.markdown(f"**{h['role'].upper()}** <span class='small-muted'>{h['ts']}</span>", unsafe_allow_html=True)
            st.write(h['content'][:200] + ('...' if len(h['content'])>200 else ''))
    else:
        st.write("No activity yet.")
    st.markdown("</div>", unsafe_allow_html=True)

st.markdown("<div style='text-align:center;margin-top:10px' class='small-muted'>No Chroma warnings — PersistentClient active · Off-white theme</div>", unsafe_allow_html=True)
